# import subprocess
# import tempfile
# import os
# import shutil
# import logging

# from core.config import get_soffice_path

# logger = logging.getLogger(__name__)

# def pdf_to_ppt(input_file):
#     soffice = get_soffice_path()
#     with tempfile.TemporaryDirectory() as tmpdir:
#         pdf_path = os.path.join(tmpdir, "input.pdf")
#         pptx_path = os.path.join(tmpdir, "input.pptx")

#         # Write uploaded PDF to temp file
#         with open(pdf_path, "wb") as f:
#             f.write(input_file.read())

#         # LibreOffice headless conversion
#         command = [
#             soffice,
#             "--headless",
#             "--convert-to",
#             "pptx",
#             pdf_path,
#             "--outdir",
#             tmpdir
#         ]

#         try:
#             result = subprocess.run(
#                 command,
#                 stdout=subprocess.PIPE,
#                 stderr=subprocess.PIPE,
#                 check=True,
#                 text=True
#             )
            
#             # If subprocess succeeds, check result output
#             logger.info("LibreOffice stdout: %s", result.stdout)
#             logger.info("LibreOffice stderr: %s", result.stderr)
        
#         except subprocess.CalledProcessError as e:
#             logger.error("LibreOffice failed: %s", e.stderr)
#             raise RuntimeError(f"LibreOffice conversion failed: {e.stderr}")

#         # Load result into memory
#         output = tempfile.SpooledTemporaryFile()
#         with open(pptx_path, "rb") as f:
#             shutil.copyfileobj(f, output)

#         output.seek(0)
#         return output





import fitz
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches
import shutil

def pdf_to_ppt(input_file):
    prs = Presentation()

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_bytes = input_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")

        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=200)
            img_path = os.path.join(tmpdir, f"page_{i}.png")
            pix.save(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                img_path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

        output = tempfile.SpooledTemporaryFile()
        prs.save(output)
        output.seek(0)
        return output